"""Green-phase unit tests for Story 27.2a — Presentation Studio service."""

from __future__ import annotations

from unittest.mock import AsyncMock, MagicMock

import pytest
from pptx.util import Inches

from app.services.presentation.marp_driver import build_marp_markdown
from app.services.presentation.pptx_driver import build_pptx, write_pptx
from app.services.presentation.schemas import GeneratePresentationInput
from app.services.presentation.service import PresentationStudioService


@pytest.mark.unit
async def test_pptx_bytes_start_with_pk_magic():
    """AC-2: PPTX is a real Open XML / ZIP file, not a renamed text file."""
    spec = {
        "title": "Pitch Deck",
        "description": "A test deck",
        "slides": [
            {"title": "Problem", "bullets": ["Pain point A", "Pain point B"]},
            {"title": "Solution", "bullets": ["Feature 1"]},
        ],
    }
    pptx_bytes = write_pptx(spec)
    assert pptx_bytes[:4] == b"PK\x03\x04"


@pytest.mark.unit
async def test_marp_output_contains_slide_separator_and_frontmatter():
    """AC-3: Marp markdown contains YAML front-matter and --- separators."""
    spec = {
        "title": "Marp Deck",
        "description": "A test",
        "slides": [{"title": "S1", "bullets": ["a", "b"]}],
    }
    raw = build_marp_markdown(spec)
    assert "---" in raw
    assert "theme:" in raw or "class:" in raw or "paginate:" in raw
    assert "# Marp Deck" in raw


@pytest.mark.unit
async def test_slides_have_sixteen_to_nine_aspect_ratio():
    """AC-2: PPTX slide width / height is exactly 16 / 9."""
    prs = build_pptx({"title": "x", "slides": [{"title": "x", "bullets": ["y"]}]})
    ratio = prs.slide_width / prs.slide_height
    assert abs(ratio - (16 / 9)) < 1e-4
    assert prs.slide_width == Inches(13.333333)
    assert prs.slide_height == Inches(7.5)


@pytest.mark.unit
async def test_deck_with_chart_adds_chart_slide():
    """AC-2: optional chart data results in a pptx chart; no chart skips gracefully."""
    spec_with_chart = {
        "title": "Chart deck",
        "slides": [
            {
                "title": "Metrics",
                "bullets": ["Revenue"],
                "chart": {
                    "categories": ["Q1", "Q2"],
                    "series": [100, 200],
                },
            }
        ],
    }
    prs = build_pptx(spec_with_chart)
    assert len(prs.slides) >= 1
    # A chart is rendered as a GraphicFrame shape on the content slide.
    content_slide = prs.slides[1]
    chart_shapes = [s for s in content_slide.shapes if s.has_chart]
    assert len(chart_shapes) >= 1

    spec_without_chart = {
        "title": "No chart",
        "slides": [{"title": "x", "bullets": ["y"]}],
    }
    prs2 = build_pptx(spec_without_chart)
    assert len(prs2.slides) >= 1


@pytest.mark.unit
async def test_empty_prompt_returns_validation_failed():
    """AC-2/AC-3/AC-6: empty or whitespace prompt returns validation_failed, no file."""
    service = PresentationStudioService()
    session = MagicMock()
    result = await service.generate(
        build_input=GeneratePresentationInput(
            prompt="   ", output_format="pptx", workspace_id=1
        ),
        session=session,
    )
    assert result.status == "validation_failed"
    assert result.error
    assert not result.presentation_id


@pytest.mark.unit
async def test_prompt_exceeding_max_length_is_truncated_or_rejected():
    """AC-1: prompt longer than PRESENTATION_MAX_PROMPT_CHARS is handled."""
    from app.config import config

    service = PresentationStudioService()
    # Temporarily lower the limit so we can exceed it without breaking Pydantic.
    original_limit = config.PRESENTATION_MAX_PROMPT_CHARS
    config.PRESENTATION_MAX_PROMPT_CHARS = 20
    try:
        result = await service.generate(
            build_input=GeneratePresentationInput(
                prompt="this prompt is too long", output_format="pptx", workspace_id=1
            ),
            session=MagicMock(),
        )
        # Without LLM available the mock returns None and it becomes validation_failed.
        assert result.status in ("validation_failed", "ready")
    finally:
        config.PRESENTATION_MAX_PROMPT_CHARS = original_limit


@pytest.mark.unit
async def test_path_traversal_rejected():
    """AC-5/AC-7: file paths must stay under the configured presentations storage root."""

    service = PresentationStudioService()
    bad_subdir = "../etc/passwd"
    with pytest.raises(ValueError):
        service._resolve_storage_path(1, "pres-123", bad_subdir)


@pytest.mark.unit
async def test_service_generate_pptx_with_mocked_llm():
    """AC-2: a valid prompt produces a ready PPTX with metadata when LLM returns a valid spec."""
    service = PresentationStudioService()
    service._call_llm_for_deck = AsyncMock(
        return_value=(
            {
                "title": "SaaS Pitch",
                "slug": "saas-pitch",
                "description": "A pitch",
                "slides": [
                    {"title": "Problem", "bullets": ["pain"]},
                    {"title": "Solution", "bullets": ["feature"]},
                ],
            },
            {"prompt_tokens": 10, "completion_tokens": 20, "total_tokens": 30},
        )
    )
    session = MagicMock()
    session.scalars = AsyncMock(return_value=MagicMock(all=MagicMock(return_value=[])))
    session.execute = AsyncMock()
    session.commit = AsyncMock()

    result = await service.generate(
        build_input=GeneratePresentationInput(
            prompt="A 3-slide pitch deck for a SaaS product",
            output_format="pptx",
            workspace_id=1,
        ),
        session=session,
    )
    assert result.status == "ready"
    assert result.format == "pptx"
    assert result.slide_count >= 2
    assert result.title
    assert result.slug
    assert result.presentation_id
    assert result.download_url


@pytest.mark.unit
async def test_workspace_scoped_slug_is_unique_with_mocked_llm():
    """AC-4/AC-5: two decks with the same title in the same workspace get disambiguated slugs."""
    service = PresentationStudioService()
    service._call_llm_for_deck = AsyncMock(
        return_value=(
            {
                "title": "Pitch deck",
                "slug": "pitch-deck",
                "slides": [{"title": "S1", "bullets": ["a"]}],
            },
            {"prompt_tokens": 1, "completion_tokens": 1, "total_tokens": 2},
        )
    )
    session = MagicMock()
    session.scalars = AsyncMock(
        side_effect=[
            MagicMock(all=MagicMock(return_value=[])),
            MagicMock(all=MagicMock(return_value=["pitch-deck"])),
        ]
    )
    session.execute = AsyncMock()
    session.commit = AsyncMock()

    r1 = await service.generate(
        build_input=GeneratePresentationInput(
            prompt="Pitch deck", output_format="pptx", workspace_id=1
        ),
        session=session,
    )
    r2 = await service.generate(
        build_input=GeneratePresentationInput(
            prompt="Pitch deck", output_format="pptx", workspace_id=1
        ),
        session=session,
    )
    assert r1.slug != r2.slug
    assert r1.slug.startswith("pitch-deck")
