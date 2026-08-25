"""PPTX writer for Presentation Studio (Story 27.2a)."""

from __future__ import annotations

import io
from typing import Any

try:
    from pptx import Presentation
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.util import Inches
except ImportError:  # pragma: no cover - exercised when python-pptx is absent
    Presentation = None  # type: ignore[assignment,misc]
    ChartData = None  # type: ignore[assignment,misc]
    XL_CHART_TYPE = None  # type: ignore[assignment,misc]
    Inches = None  # type: ignore[assignment,misc]


def _require_pptx() -> None:
    if Presentation is None or Inches is None:
        raise RuntimeError("python-pptx is not installed; cannot generate PPTX decks")


def _set_shape_text(shape: Any, text: str) -> None:
    """Best-effort text assignment when a placeholder may be missing."""
    if shape is None:
        return
    try:
        shape.text = text
    except Exception:
        return


def _set_notes(slide: Any, notes: str) -> None:
    try:
        slide.notes_slide.notes_text_frame.text = notes
    except Exception:
        return


def build_pptx(deck_spec: dict[str, Any]) -> Any:
    """Build a 16:9 PPTX from a DeckSpec dict.

    The deck must contain at least a title and one or more slide dicts.
    Charts are added only when ``chart`` has numeric ``series``.
    """
    _require_pptx()
    assert Presentation is not None and Inches is not None

    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    # Title slide layout
    title_slide_layout = prs.slide_layouts[0]
    title_slide = prs.slides.add_slide(title_slide_layout)
    title = deck_spec.get("title", "Untitled")
    if title_slide.shapes.title is not None:
        _set_shape_text(title_slide.shapes.title, title)
    if len(title_slide.placeholders) > 1:
        _set_shape_text(title_slide.placeholders[1], deck_spec.get("description") or "")
    if description := deck_spec.get("description"):
        _set_notes(title_slide, str(description))

    # Content layout
    content_layout = prs.slide_layouts[1]
    for slide_spec in deck_spec.get("slides", []):
        slide = prs.slides.add_slide(content_layout)
        if slide.shapes.title is not None:
            _set_shape_text(slide.shapes.title, slide_spec.get("title", ""))
        if len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for i, bullet in enumerate(slide_spec.get("bullets", [])):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = str(bullet)
                p.level = 0

        # Optional chart
        chart_data = slide_spec.get("chart")
        if (
            ChartData is not None
            and XL_CHART_TYPE is not None
            and chart_data
            and chart_data.get("series")
        ):
            categories = [str(c) for c in chart_data.get("categories", [])]
            values: list[float] = []
            for v in chart_data.get("series", []):
                if v is None:
                    continue
                try:
                    values.append(float(v))
                except (TypeError, ValueError):
                    continue
            n = min(len(categories), len(values))
            categories, values = categories[:n], values[:n]
            if categories and values:
                chart = ChartData()
                chart.title = slide_spec.get("title", "")
                chart.categories = categories
                chart.add_series("Series 1", values)
                x, y, cx, cy = Inches(7.5), Inches(1.5), Inches(5.0), Inches(4.5)
                slide.shapes.add_chart(
                    XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart
                )

        if notes := slide_spec.get("notes"):
            _set_notes(slide, str(notes))

    return prs


def write_pptx(deck_spec: dict[str, Any]) -> bytes:
    """Return PPTX bytes for the given deck spec."""
    prs = build_pptx(deck_spec)
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
